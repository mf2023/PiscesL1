#!/usr/bin/env/python3
# -*- coding: utf-8 -*-

# Copyright © 2025-2026 Wenze Wei. All Rights Reserved.
#
# This file is part of PiscesL1.
# The PiscesL1 project belongs to the Dunimd Team.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# You may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
Document Processor Tools - PDF, Word, PowerPoint processing
"""

import os
from typing import Any, Dict, List, Optional

from .base import POPSSMCPToolBase, POPSSMCPToolResult


class PDFReaderTool(POPSSMCPToolBase):
    name = "pdf_read"
    description = "Extract text content from PDF files"
    category = "document"
    tags = ["pdf", "document", "extract", "text"]
    
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to PDF file"
            },
            "start_page": {
                "type": "integer",
                "description": "Start page number (1-indexed)",
                "default": 1
            },
            "end_page": {
                "type": "integer",
                "description": "End page number (inclusive)",
                "default": None
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to extract",
                "default": 50000
            }
        },
        "required": ["file_path"]
    }
    
    async def execute(self, arguments: Dict[str, Any]) -> POPSSMCPToolResult:
        file_path = arguments.get("file_path", "")
        start_page = arguments.get("start_page", 1)
        end_page = arguments.get("end_page")
        max_chars = arguments.get("max_chars", 50000)
        
        if not file_path:
            return self._create_error_result("file_path is required", "ValidationError")
        
        if not os.path.exists(file_path):
            return self._create_error_result(f"File not found: {file_path}", "FileNotFoundError")
        
        try:
            import fitz
            
            doc = fitz.open(file_path)
            total_pages = len(doc)
            
            start = max(0, start_page - 1)
            end = min(total_pages, end_page) if end_page else total_pages
            
            text_parts = []
            char_count = 0
            
            for page_num in range(start, end):
                page = doc[page_num]
                text = page.get_text()
                
                if char_count + len(text) > max_chars:
                    text = text[:max_chars - char_count]
                    text_parts.append(text)
                    break
                
                text_parts.append(text)
                char_count += len(text)
            
            doc.close()
            
            full_text = "\n".join(text_parts)
            
            return self._create_success_result({
                "file_path": file_path,
                "total_pages": total_pages,
                "extracted_pages": f"{start + 1}-{end}",
                "char_count": len(full_text),
                "text": full_text,
            })
            
        except ImportError:
            return self._create_error_result(
                "PyMuPDF (fitz) required. Install with: pip install PyMuPDF",
                "DependencyError"
            )
        except Exception as e:
            return self._create_error_result(str(e), type(e).__name__)


class WordReaderTool(POPSSMCPToolBase):
    name = "word_read"
    description = "Extract text content from Word documents (.docx)"
    category = "document"
    tags = ["word", "docx", "document", "extract"]
    
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to Word document"
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to extract",
                "default": 50000
            }
        },
        "required": ["file_path"]
    }
    
    async def execute(self, arguments: Dict[str, Any]) -> POPSSMCPToolResult:
        file_path = arguments.get("file_path", "")
        max_chars = arguments.get("max_chars", 50000)
        
        if not file_path:
            return self._create_error_result("file_path is required", "ValidationError")
        
        if not os.path.exists(file_path):
            return self._create_error_result(f"File not found: {file_path}", "FileNotFoundError")
        
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            paragraphs = []
            char_count = 0
            
            for para in doc.paragraphs:
                text = para.text
                if char_count + len(text) > max_chars:
                    text = text[:max_chars - char_count]
                    paragraphs.append(text)
                    break
                paragraphs.append(text)
                char_count += len(text)
            
            full_text = "\n".join(paragraphs)
            
            return self._create_success_result({
                "file_path": file_path,
                "paragraph_count": len(doc.paragraphs),
                "char_count": len(full_text),
                "text": full_text,
            })
            
        except ImportError:
            return self._create_error_result(
                "python-docx required. Install with: pip install python-docx",
                "DependencyError"
            )
        except Exception as e:
            return self._create_error_result(str(e), type(e).__name__)


class PowerPointReaderTool(POPSSMCPToolBase):
    name = "ppt_read"
    description = "Extract text content from PowerPoint presentations (.pptx)"
    category = "document"
    tags = ["powerpoint", "pptx", "document", "extract"]
    
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to PowerPoint file"
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to extract",
                "default": 50000
            }
        },
        "required": ["file_path"]
    }
    
    async def execute(self, arguments: Dict[str, Any]) -> POPSSMCPToolResult:
        file_path = arguments.get("file_path", "")
        max_chars = arguments.get("max_chars", 50000)
        
        if not file_path:
            return self._create_error_result("file_path is required", "ValidationError")
        
        if not os.path.exists(file_path):
            return self._create_error_result(f"File not found: {file_path}", "FileNotFoundError")
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_content = []
            char_count = 0
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                
                slide_content = f"--- Slide {i} ---\n" + "\n".join(slide_text)
                
                if char_count + len(slide_content) > max_chars:
                    break
                
                slides_content.append(slide_content)
                char_count += len(slide_content)
            
            full_text = "\n\n".join(slides_content)
            
            return self._create_success_result({
                "file_path": file_path,
                "slide_count": len(prs.slides),
                "extracted_slides": len(slides_content),
                "char_count": len(full_text),
                "text": full_text,
            })
            
        except ImportError:
            return self._create_error_result(
                "python-pptx required. Install with: pip install python-pptx",
                "DependencyError"
            )
        except Exception as e:
            return self._create_error_result(str(e), type(e).__name__)


class TextFileReaderTool(POPSSMCPToolBase):
    name = "text_read"
    description = "Read text files with encoding detection"
    category = "document"
    tags = ["text", "file", "read"]
    
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to text file"
            },
            "encoding": {
                "type": "string",
                "description": "File encoding (default: auto-detect)",
                "default": None
            },
            "max_chars": {
                "type": "integer",
                "description": "Maximum characters to read",
                "default": 100000
            }
        },
        "required": ["file_path"]
    }
    
    async def execute(self, arguments: Dict[str, Any]) -> POPSSMCPToolResult:
        file_path = arguments.get("file_path", "")
        encoding = arguments.get("encoding")
        max_chars = arguments.get("max_chars", 100000)
        
        if not file_path:
            return self._create_error_result("file_path is required", "ValidationError")
        
        if not os.path.exists(file_path):
            return self._create_error_result(f"File not found: {file_path}", "FileNotFoundError")
        
        try:
            if encoding:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read(max_chars)
            else:
                encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
                content = None
                used_encoding = None
                
                for enc in encodings:
                    try:
                        with open(file_path, "r", encoding=enc) as f:
                            content = f.read(max_chars)
                        used_encoding = enc
                        break
                    except UnicodeDecodeError:
                        continue
                
                if content is None:
                    with open(file_path, "rb") as f:
                        content = f.read(max_chars).decode("utf-8", errors="replace")
                    used_encoding = "binary"
            
            return self._create_success_result({
                "file_path": file_path,
                "encoding": used_encoding or encoding,
                "char_count": len(content),
                "text": content,
            })
            
        except Exception as e:
            return self._create_error_result(str(e), type(e).__name__)
