#!/usr/bin/env/python3

# Copyright © 2025 Wenze Wei. All Rights Reserved.
#
# This file is part of PiscesL1.
# The PiscesL1 project belongs to the Dunimd project team.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# You may not use this file except in compliance with the License.
# Commercial use is strictly prohibited.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific governing permissions and
# limitations under the License.

import os
import sys
import json
import fitz
import docx
import pptx
from pathlib import Path
from pathlib import Path
from pptx import Presentation
from typing import Dict, Any, List, Optional

# Add project root to path
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

from utils.mcp import PiscesLxCoreMCPPlaza

# Create mcp instance for tool registration
mcp = PiscesLxCoreMCPPlaza()

class DocumentProcessor:
    """
    Document processor for PDF, DOCX, and PPTX files with security constraints.
    
    This class provides methods to extract text, metadata, and structural information
    from various document formats while enforcing security limits to prevent
    resource exhaustion and path traversal vulnerabilities.
    """
    
    def __init__(self):
        """
        Initialize DocumentProcessor with supported formats and security limits.
        """
        self.supported_formats = {
            '.pdf': 'PDF Document',
            '.docx': 'Word Document',
            '.doc': 'Word 97-2003 Document',
            '.pptx': 'PowerPoint Presentation',
            '.ppt': 'PowerPoint 97-2003 Presentation'
        }
        # Security limits for server environment
        self.max_file_size = 50 * 1024 * 1024  # 50MB limit
        self.max_pages = 100  # Maximum pages per document
        self.max_extract_chars = 100000  # Maximum characters to extract per document
        self.allowed_mime_types = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.doc': 'application/msword',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.ppt': 'application/vnd.ms-powerpoint'
        }
    
    def _validate_file(self, file_path: str) -> Path:
        """
        Validate file path and check if it exists with security checks.
        
        Args:
            file_path (str): Path to the file to validate
            
        Returns:
            Path: Validated Path object
            
        Raises:
            FileNotFoundError: If file does not exist
            ValueError: If file format is unsupported or security checks fail
        """
        path = Path(file_path)
        
        # Check if file exists
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Check file format support
        if path.suffix.lower() not in self.supported_formats:
            raise ValueError(f"Unsupported format: {path.suffix}")
        
        # Security checks
        file_size = path.stat().st_size
        if file_size > self.max_file_size:
            raise ValueError(f"File too large: {file_size} bytes (max: {self.max_file_size})")
        
        # Additional security: check for potential path traversal
        try:
            resolved_path = path.resolve()
            if not str(resolved_path).startswith(str(Path.cwd())):
                raise ValueError("File path outside allowed directory")
        except Exception:
            raise ValueError("Invalid file path")
        
        return path
    
    def extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content from PDF file with security limits.
        
        Args:
            file_path (str): Path to the PDF file
            
        Returns:
            Dict[str, Any]: Dictionary containing extracted content or error information
        """
        try:
            path = self._validate_file(file_path)
            doc = fitz.open(str(path))
            
            # Security check: limit pages
            if len(doc) > self.max_pages:
                doc.close()
                raise ValueError(f"PDF too long: {len(doc)} pages (max: {self.max_pages})")
            
            content = {
                "pages": [],
                "metadata": doc.metadata,
                "total_pages": len(doc),
                "file_size": path.stat().st_size
            }
            
            total_chars = 0
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text()
                
                # Security check: limit extracted text
                total_chars += len(text)
                if total_chars > self.max_extract_chars:
                    doc.close()
                    raise ValueError(f"Text content too large (max: {self.max_extract_chars} chars)")
                
                # Extract tables and images (simplified for security)
                blocks = page.get_text("dict")
                tables = self._extract_tables_from_page(blocks)
                images = self._extract_images_from_page(page)
                
                content["pages"].append({
                    "page_number": page_num + 1,
                    "text": text[:10000],  # Limit text per page
                    "tables": tables[:5],   # Limit tables per page
                    "images": images[:10],  # Limit images per page
                    "word_count": len(text.split())
                })
            
            doc.close()
            return {
                "success": True,
                "format": "pdf",
                "content": content,
                "summary": {
                    "total_words": sum(p["word_count"] for p in content["pages"]),
                    "total_tables": sum(len(p["tables"]) for p in content["pages"]),
                    "total_images": sum(len(p["images"]) for p in content["pages"])
                }
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def extract_docx_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content from DOCX file with security limits.
        
        Args:
            file_path (str): Path to the DOCX file
            
        Returns:
            Dict[str, Any]: Dictionary containing extracted content or error information
        """
        try:
            path = self._validate_file(file_path)
            doc = docx.Document(str(path))
            
            content = {
                "paragraphs": [],
                "tables": [],
                "metadata": {
                    "core_properties": {
                        "title": doc.core_properties.title,
                        "author": doc.core_properties.author,
                        "subject": doc.core_properties.subject,
                        "created": str(doc.core_properties.created),
                        "modified": str(doc.core_properties.modified)
                    }
                },
                "file_size": path.stat().st_size
            }
            
            # Security: limit paragraphs and tables
            max_paragraphs = 1000
            max_tables = 50
            
            # Extract paragraphs (with limit)
            for i, para in enumerate(doc.paragraphs):
                if i >= max_paragraphs:
                    break
                if para.text.strip():  # Only include non-empty paragraphs
                    content["paragraphs"].append({
                        "text": para.text[:5000],  # Limit text length per paragraph
                        "style": para.style.name if para.style else "Normal",
                        "alignment": str(para.alignment) if para.alignment else None
                    })
            
            # Extract tables (with limit)
            for table_idx, table in enumerate(doc.tables):
                if table_idx >= max_tables:
                    break
                    
                table_data = {
                    "table_index": table_idx,
                    "rows": [],
                    "column_count": len(table.columns),
                    "row_count": len(table.rows)
                }
                
                # Limit rows per table
                max_rows = 100
                for i, row in enumerate(table.rows):
                    if i >= max_rows:
                        break
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text[:1000])  # Limit cell text
                    table_data["rows"].append(row_data)
                
                content["tables"].append(table_data)
            
            return {
                "success": True,
                "format": "docx",
                "content": content,
                "summary": {
                    "total_paragraphs": len(content["paragraphs"]),
                    "total_tables": len(content["tables"]),
                    "total_words": sum(len(p["text"].split()) for p in content["paragraphs"])
                }
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def extract_pptx_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content from PPTX file with security limits.
        
        Args:
            file_path (str): Path to the PPTX file
            
        Returns:
            Dict[str, Any]: Dictionary containing extracted content or error information
        """
        try:
            path = self._validate_file(file_path)
            prs = Presentation(str(path))
            
            # Security check: limit slides
            if len(prs.slides) > self.max_pages:
                raise ValueError(f"Presentation too long: {len(prs.slides)} slides (max: {self.max_pages})")
            
            content = {
                "slides": [],
                "metadata": {
                    "slide_width": prs.slide_width,
                    "slide_height": prs.slide_height,
                    "slide_count": len(prs.slides)
                },
                "file_size": path.stat().st_size
            }
            
            # Security: limit slides processed
            max_slides = min(len(prs.slides), self.max_pages)
            
            for slide_idx in range(max_slides):
                slide = prs.slides[slide_idx]
                slide_content = {
                    "slide_number": slide_idx + 1,
                    "title": "",
                    "text_content": [],
                    "tables": [],
                    "notes": slide.notes_slide.notes_text_frame.text[:1000] if slide.has_notes_slide else ""
                }
                
                # Extract text from shapes (with limits)
                text_count = 0
                max_text_items = 50
                max_shapes = 100
                
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape_idx >= max_shapes:
                        break
                        
                    if hasattr(shape, "text") and shape.text and text_count < max_text_items:
                        if shape.text_frame and shape.text_frame.paragraphs:
                            first_para = shape.text_frame.paragraphs[0]
                            if first_para.level == 0 and not slide_content["title"]:
                                slide_content["title"] = shape.text[:200]  # Limit title length
                            else:
                                slide_content["text_content"].append({
                                    "text": shape.text[:1000],  # Limit text length
                                    "level": first_para.level if first_para else 0
                                })
                            text_count += 1
                    
                    # Extract tables (with limits)
                    if shape.has_table and len(slide_content["tables"]) < 10:
                        table = shape.table
                        table_data = {
                            "rows": [],
                            "column_count": len(table.columns),
                            "row_count": len(table.rows)
                        }
                        
                        # Limit rows per table
                        max_table_rows = 50
                        for i, row in enumerate(table.rows):
                            if i >= max_table_rows:
                                break
                            row_data = [cell.text[:500] for cell in row.cells]  # Limit cell text
                            table_data["rows"].append(row_data)
                        
                        slide_content["tables"].append(table_data)
                
                content["slides"].append(slide_content)
            
            return {
                "success": True,
                "format": "pptx",
                "content": content,
                "summary": {
                    "total_slides": len(content["slides"]),
                    "total_tables": sum(len(s["tables"]) for s in content["slides"]),
                    "total_words": sum(
                        len(text["text"].split()) 
                        for slide in content["slides"] 
                        for text in slide["text_content"]
                    )
                }
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _extract_tables_from_page(self, blocks: Dict[str, Any]) -> List[Dict[str, Any]]:
        """
        Extract tables from PDF page blocks.
        
        Args:
            blocks (Dict[str, Any]): PDF page blocks data
            
        Returns:
            List[Dict[str, Any]]: List of extracted tables
        """
        tables = []
        # Simplified table extraction - can be enhanced
        return tables
    
    def _extract_images_from_page(self, page) -> List[Dict[str, Any]]:
        """
        Extract images from PDF page.
        
        Args:
            page: PDF page object
            
        Returns:
            List[Dict[str, Any]]: List of extracted images information
        """
        images = []
        image_list = page.get_images()
        
        for img_idx, img in enumerate(image_list):
            images.append({
                "image_index": img_idx,
                "width": img[2],
                "height": img[3],
                "colorspace": img[5],
                "size": len(img[7]) if len(img) > 7 else 0
            })
        
        return images
    
    def get_document_summary(self, file_path: str) -> Dict[str, Any]:
        """
        Get quick summary of document without full content extraction.
        
        Args:
            file_path (str): Path to the document file
            
        Returns:
            Dict[str, Any]: Dictionary containing document summary or error information
        """
        try:
            path = self._validate_file(file_path)
            extension = path.suffix.lower()
            
            if extension == '.pdf':
                doc = fitz.open(str(path))
                summary = {
                    "format": "pdf",
                    "pages": len(doc),
                    "metadata": doc.metadata,
                    "file_size": path.stat().st_size
                }
                doc.close()
            elif extension == '.docx':
                doc = docx.Document(str(path))
                summary = {
                    "format": "docx",
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables),
                    "metadata": {
                        "title": doc.core_properties.title,
                        "author": doc.core_properties.author
                    },
                    "file_size": path.stat().st_size
                }
            elif extension == '.pptx':
                prs = Presentation(str(path))
                summary = {
                    "format": "pptx",
                    "slides": len(prs.slides),
                    "file_size": path.stat().st_size
                }
            else:
                return {"success": False, "error": f"Unsupported format: {extension}"}
            
            return {"success": True, "summary": summary}
            
        except Exception as e:
            return {"success": False, "error": str(e)}

# Initialize document processor
document_processor = DocumentProcessor()

@mcp.tool()
def extract_document_content(file_path: str, include_full_content: bool = True) -> Dict[str, Any]:
    """
    Extract content from PDF, DOCX, or PPTX files with security restrictions.
    
    Args:
        file_path (str): Path to the document file (max 50MB, 100 pages/slides)
        include_full_content (bool): Whether to include full content or just summary
        
    Returns:
        Dict[str, Any]: Dictionary containing extracted content and metadata
    """
    if not include_full_content:
        return document_processor.get_document_summary(file_path)
    
    path = Path(file_path)
    extension = path.suffix.lower()
    
    if extension == '.pdf':
        return document_processor.extract_pdf_content(file_path)
    elif extension in ['.docx', '.doc']:
        return document_processor.extract_docx_content(file_path)
    elif extension in ['.pptx', '.ppt']:
        return document_processor.extract_pptx_content(file_path)
    else:
        return {
            "success": False,
            "error": f"Unsupported format: {extension}",
            "supported_formats": list(document_processor.supported_formats.keys())
        }

@mcp.tool()
def list_supported_formats() -> Dict[str, Any]:
    """
    List all supported document formats.
    
    Returns:
        Dict[str, Any]: Dictionary containing supported formats
    """
    return {
        "success": True,
        "formats": document_processor.supported_formats,
        "description": "Supported document formats for processing"
    }

@mcp.tool()
def batch_process_documents(directory_path: str, recursive: bool = False) -> Dict[str, Any]:
    """
    Process multiple documents in a directory with security restrictions.
    
    Args:
        directory_path (str): Path to the directory containing documents
        recursive (bool): Whether to search subdirectories
        
    Returns:
        Dict[str, Any]: List of processed documents with summaries
    """
    try:
        path = Path(directory_path)
        if not path.exists() or not path.is_dir():
            return {"success": False, "error": "Directory not found"}
        
        # Security: resolve and validate directory path
        resolved_path = path.resolve()
        if not str(resolved_path).startswith(str(Path.cwd())):
            return {"success": False, "error": "Directory path outside allowed location"}
        
        documents = []
        pattern = "**/*" if recursive else "*"
        
        # Security: limit number of files processed
        max_documents = 100
        file_count = 0
        
        for file_path in path.glob(pattern):
            if file_count >= max_documents:
                documents.append({
                    "file_path": "LIMIT_REACHED",
                    "summary": {"error": "Maximum document limit reached (100)"}
                })
                break
                
            if file_path.suffix.lower() in document_processor.supported_formats:
                result = document_processor.get_document_summary(str(file_path))
                if result["success"]:
                    documents.append({
                        "file_path": str(file_path),
                        "summary": result["summary"]
                    })
                    file_count += 1
        
        return {
            "success": True,
            "directory": str(path),
            "documents": documents,
            "count": len(documents),
            "security_notice": "Limited to 100 documents per batch for server safety"
        }
        
    except Exception as e:
        return {"success": False, "error": str(e)}
